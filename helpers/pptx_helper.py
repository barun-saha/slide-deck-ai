"""
A set of functions to create a PowerPoint slide deck.
"""
import logging
import pathlib
import random
import re
import sys
import tempfile
from typing import List, Tuple, Optional

import json5
import pptx
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

sys.path.append('..')
sys.path.append('../..')

import helpers.image_search as ims
from global_config import GlobalConfig


load_dotenv()


# English Metric Unit (used by PowerPoint) to inches
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)

STEP_BY_STEP_PROCESS_MARKER = '>> '
IMAGE_DISPLAY_PROBABILITY = 0.3
FOREGROUND_IMAGE_PROBABILITY = 0.75
SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)


logger = logging.getLogger(__name__)


def remove_slide_number_from_heading(header: str) -> str:
    """
    Remove the slide number from a given slide header.

    :param header: The header of a slide.
    :return: The header without slide number.
    """

    if SLIDE_NUMBER_REGEX.match(header):
        idx = header.find(':')
        header = header[idx + 1:]

    return header


def generate_powerpoint_presentation(
        structured_data: str,
        slides_template: str,
        output_file_path: pathlib.Path
) -> List:
    """
    Create and save a PowerPoint presentation file containing the content in JSON format.

    :param structured_data: The presentation contents as "JSON" (may contain trailing commas).
    :param slides_template: The PPTX template to use.
    :param output_file_path: The path of the PPTX file to save as.
    :return: A list of presentation title and slides headers.
    """

    # The structured "JSON" might contain trailing commas, so using json5
    parsed_data = json5.loads(structured_data)
    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])
    slide_width_inch, slide_height_inch = _get_slide_width_height_inches(presentation)

    # The title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = parsed_data['title']
    logger.info(
        'PPT title: %s | #slides: %d | template: %s',
        title.text, len(parsed_data['slides']),
        GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file']
    )
    subtitle.text = 'by Myself and SlideDeck AI :)'
    all_headers = [title.text, ]

    # Add content in a loop
    for a_slide in parsed_data['slides']:
        is_processing_done = _handle_double_col_layout(
            presentation=presentation,
            slide_json=a_slide,
            slide_width_inch=slide_width_inch,
            slide_height_inch=slide_height_inch
        )

        if not is_processing_done:
            is_processing_done = _handle_step_by_step_process(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

        if not is_processing_done:
            _handle_default_display(
                presentation=presentation,
                slide_json=a_slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch
            )

    # The thank-you slide
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    title = slide.shapes.title
    title.text = 'Thank you!'

    presentation.save(output_file_path)

    return all_headers


def get_flat_list_of_contents(items: list, level: int) -> List[Tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and
    its level.

    :param items: A bullet point (string or list).
    :param level: The current level of hierarchy.
    :return: A list of (bullet item text, hierarchical level) tuples.
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list


def get_slide_placeholders(
        slide: pptx.slide.Slide,
        layout_number: int,
        is_debug: bool = False
) -> List[Tuple[int, str]]:
    """
    Return the index and name (lower case) of all placeholders present in a slide, except
    the title placeholder.

    A placeholder in a slide is a place to add content. Each placeholder has a name and an index.
    This index is NOT a list index, rather a set of keys used to look up a dict. So, `idx` is
    non-contiguous. Also, the title placeholder of a slide always has index 0. User-added
    placeholder get indices assigned starting from 10.

    With user-edited or added placeholders, their index may be difficult to track. This function
    returns the placeholders name as well, which could be useful to distinguish between the
    different placeholder.

    :param slide: The slide.
    :param layout_number: The layout number used by the slide.
    :param is_debug: Whether to print debugging statements.
    :return: A list containing placeholders (idx, name) tuples, except the title placeholder.
    """

    if is_debug:
        print(
            f'Slide layout #{layout_number}:'
            f' # of placeholders: {len(slide.shapes.placeholders)} (including the title)'
        )

    placeholders = [
        (shape.placeholder_format.idx, shape.name.lower()) for shape in slide.shapes.placeholders
    ]
    placeholders.pop(0)  # Remove the title placeholder

    if is_debug:
        print(placeholders)

    return placeholders


def _handle_default_display(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Display a list of text in a slide.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    """

    status = False

    if random.random() < IMAGE_DISPLAY_PROBABILITY:
        if random.random() < FOREGROUND_IMAGE_PROBABILITY:
            status = _handle_display_image__in_foreground(
                presentation,
                slide_json,
                slide_width_inch,
                slide_height_inch
            )
        else:
            status = _handle_display_image__in_background(
                presentation,
                slide_json,
                slide_width_inch,
                slide_height_inch
            )

    if status:
        return

    # Image display failed, so display only text
    bullet_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(bullet_slide_layout)

    shapes = slide.shapes
    title_shape = shapes.title

    try:
        body_shape = shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        body_shape = shapes.placeholders[placeholders[0][0]]

    title_shape.text = remove_slide_number_from_heading(slide_json['heading'])
    text_frame = body_shape.text_frame

    # The bullet_points may contain a nested hierarchy of JSON arrays
    # In some scenarios, it may contain objects (dictionaries) because the LLM generated so
    #  ^ The second scenario is not covered

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    _handle_key_message(
        the_slide=slide,
        slide_json=slide_json,
        slide_height_inch=slide_height_inch,
        slide_width_inch=slide_width_inch
    )


def _handle_display_image__in_foreground(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Create a slide with text and image using a picture placeholder layout.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if the side has been processed.
    """

    img_keywords = slide_json['img_keywords'].strip()
    slide = presentation.slide_layouts[8]  # Picture with Caption
    slide = presentation.slides.add_slide(slide)
    placeholders = None

    title_placeholder = slide.shapes.title
    title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])

    try:
        pic_col: PicturePlaceholder = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=8)
        pic_col = None
        for idx, name in placeholders:
            if 'picture' in name:
                pic_col: PicturePlaceholder = slide.shapes.placeholders[idx]

    try:
        text_col: SlidePlaceholder = slide.shapes.placeholders[2]
    except KeyError:
        text_col = None
        if not placeholders:
            placeholders = get_slide_placeholders(slide, layout_number=8)

        for idx, name in placeholders:
            if 'content' in name:
                text_col: SlidePlaceholder = slide.shapes.placeholders[idx]

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            text_col.text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = text_col.text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    if not img_keywords:
        # No keywords, so no image search and addition
        return True

    try:
        photo_url, page_url = ims.get_photo_url_from_api_response(
            ims.search_pexels(query=img_keywords, size='medium')
        )

        if photo_url:
            pic_col.insert_picture(
                ims.get_image_from_url(photo_url)
            )

            _add_text_at_bottom(
                slide=slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch,
                text='Photo provided by Pexels',
                hyperlink=page_url
            )
    except Exception as ex:
        logger.error(
            '*** Error occurred while running adding image to slide: %s',
            str(ex)
        )

    return True


def _handle_display_image__in_background(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with text and an image in the background. It works just like
    `_handle_default_display()` but with a background image added.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if the slide has been processed.
    """

    img_keywords = slide_json['img_keywords'].strip()

    # Add a photo in the background, text in the foreground
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title

    try:
        body_shape = slide.shapes.placeholders[1]
    except KeyError:
        placeholders = get_slide_placeholders(slide, layout_number=1)
        # Layout 1 usually has two placeholders, including the title
        body_shape = slide.shapes.placeholders[placeholders[0][0]]

    title_shape.text = remove_slide_number_from_heading(slide_json['heading'])

    flat_items_list = get_flat_list_of_contents(slide_json['bullet_points'], level=0)

    for idx, an_item in enumerate(flat_items_list):
        if idx == 0:
            body_shape.text_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        else:
            paragraph = body_shape.text_frame.add_paragraph()
            paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
            paragraph.level = an_item[1]

    if not img_keywords:
        # No keywords, so no image search and addition
        return True

    try:
        photo_url, page_url = ims.get_photo_url_from_api_response(
            ims.search_pexels(query=img_keywords, size='large')
        )

        if photo_url:
            picture = slide.shapes.add_picture(
                image_file=ims.get_image_from_url(photo_url),
                left=0,
                top=0,
                width=pptx.util.Inches(slide_width_inch),
            )

            _add_text_at_bottom(
                slide=slide,
                slide_width_inch=slide_width_inch,
                slide_height_inch=slide_height_inch,
                text='Photo provided by Pexels',
                hyperlink=page_url
            )

            # Move picture to background
            # https://github.com/scanny/python-pptx/issues/49#issuecomment-137172836
            slide.shapes._spTree.remove(picture._element)
            slide.shapes._spTree.insert(2, picture._element)
    except Exception as ex:
        logger.error(
            '*** Error occurred while running adding image to the slide background: %s',
            str(ex)
        )

    return True


def _add_text_at_bottom(
        slide: pptx.slide.Slide,
        slide_width_inch: float,
        slide_height_inch: float,
        text: str,
        hyperlink: Optional[str] = None,
        target_height: Optional[float] = 0.5
):
    """
    Add arbitrary text to a textbox positioned near the lower left side of a slide.

    :param slide: The slide.
    :param slide_width_inch: The width of the slide.
    :param slide_height_inch: The height of the slide.
    :param target_height: the target height of the box in inches (optional).
    :param text: The text to be added
    :param hyperlink: The hyperlink to be added to the text (optional).
    """

    footer = slide.shapes.add_textbox(
        left=INCHES_1,
        top=pptx.util.Inches(slide_height_inch - target_height),
        width=pptx.util.Inches(slide_width_inch),
        height=pptx.util.Inches(target_height)
    )

    paragraph = footer.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = pptx.util.Pt(10)
    run.font.underline = False

    if hyperlink:
        run.hyperlink.address = hyperlink


def _handle_double_col_layout(
        presentation: pptx.Presentation(),
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add a slide with a double column layout for comparison.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return: True if double col layout has been added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        double_col_content = slide_json['bullet_points']

        if double_col_content and (
                len(double_col_content) == 2
        ) and isinstance(double_col_content[0], dict) and isinstance(double_col_content[1], dict):
            slide = presentation.slide_layouts[4]
            slide = presentation.slides.add_slide(slide)
            placeholders = None

            shapes = slide.shapes
            title_placeholder = shapes.title
            title_placeholder.text = remove_slide_number_from_heading(slide_json['heading'])

            try:
                left_heading, right_heading = shapes.placeholders[1], shapes.placeholders[3]
            except KeyError:
                # For manually edited/added master slides, the placeholder idx numbers in the dict
                # will be different (>= 10)
                left_heading, right_heading = None, None
                placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'text placeholder' in name:
                        if not left_heading:
                            left_heading = shapes.placeholders[idx]
                        elif not right_heading:
                            right_heading = shapes.placeholders[idx]

            try:
                left_col, right_col = shapes.placeholders[2], shapes.placeholders[4]
            except KeyError:
                left_col, right_col = None, None
                if not placeholders:
                    placeholders = get_slide_placeholders(slide, layout_number=4)

                for idx, name in placeholders:
                    if 'content placeholder' in name:
                        if not left_col:
                            left_col = shapes.placeholders[idx]
                        elif not right_col:
                            right_col = shapes.placeholders[idx]

            left_col_frame, right_col_frame = left_col.text_frame, right_col.text_frame

            if 'heading' in double_col_content[0] and left_heading:
                left_heading.text = double_col_content[0]['heading']
            if 'bullet_points' in double_col_content[0]:
                flat_items_list = get_flat_list_of_contents(
                    double_col_content[0]['bullet_points'], level=0
                )

                if not left_heading:
                    left_col_frame.text = double_col_content[0]['heading']

                for idx, an_item in enumerate(flat_items_list):
                    if left_heading and idx == 0:
                        left_col_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                    else:
                        paragraph = left_col_frame.add_paragraph()
                        paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                        paragraph.level = an_item[1]

            if 'heading' in double_col_content[1] and right_heading:
                right_heading.text = double_col_content[1]['heading']
            if 'bullet_points' in double_col_content[1]:
                flat_items_list = get_flat_list_of_contents(
                    double_col_content[1]['bullet_points'], level=0
                )

                if not right_heading:
                    right_col_frame.text = double_col_content[1]['heading']

                for idx, an_item in enumerate(flat_items_list):
                    if right_col_frame and idx == 0:
                        right_col_frame.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                    else:
                        paragraph = right_col_frame.add_paragraph()
                        paragraph.text = an_item[0].removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                        paragraph.level = an_item[1]

            _handle_key_message(
                the_slide=slide,
                slide_json=slide_json,
                slide_height_inch=slide_height_inch,
                slide_width_inch=slide_width_inch
            )

            return True

    return False


def _handle_step_by_step_process(
        presentation: pptx.Presentation,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
) -> bool:
    """
    Add shapes to display a step-by-step process in the slide, if available.

    :param presentation: The presentation object.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    :return True if this slide has a step-by-step process depiction added; False otherwise.
    """

    if 'bullet_points' in slide_json and slide_json['bullet_points']:
        steps = slide_json['bullet_points']

        no_marker_count = 0.0
        n_steps = len(steps)

        # Ensure that it is a single list of strings without any sub-list
        for step in steps:
            if not isinstance(step, str):
                return False

            # In some cases, one or two steps may not begin with >>, e.g.:
            # {
            #     "heading": "Step-by-Step Process: Creating a Legacy",
            #     "bullet_points": [
            #         "Identify your unique talents and passions",
            #         ">> Develop your skills and knowledge",
            #         ">> Create meaningful work",
            #         ">> Share your work with the world",
            #         ">> Continuously learn and adapt"
            #     ],
            #     "key_message": ""
            # },
            #
            # Use a threshold, e.g., at most 20%
            if not step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                no_marker_count += 1

        slide_header = slide_json['heading'].lower()
        if (no_marker_count / n_steps > 0.25) and not (
                ('step-by-step' in slide_header) or ('step by step' in slide_header)
        ):
            return False

        if n_steps < 3 or n_steps > 6:
            # Two steps -- probably not a process
            # More than 5--6 steps -- would likely cause a visual clutter
            return False

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = remove_slide_number_from_heading(slide_json['heading'])

        if 3 <= n_steps <= 4:
            # Horizontal display
            height = INCHES_1_5
            width = pptx.util.Inches(slide_width_inch / n_steps - 0.01)
            top = pptx.util.Inches(slide_height_inch / 2)
            left = pptx.util.Inches((slide_width_inch - width.inches * n_steps) / 2 + 0.05)

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
                shape.text = step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                left += width - INCHES_0_4
        elif 4 < n_steps <= 6:
            # Vertical display
            height = pptx.util.Inches(0.65)
            top = pptx.util.Inches(slide_height_inch / 4)
            left = INCHES_1  # slide_width_inch - width.inches)

            # Find the close to median width, based on the length of each text, to be set
            # for the shapes
            width = pptx.util.Inches(slide_width_inch * 2 / 3)
            lengths = [len(step) for step in steps]
            font_size_20pt = pptx.util.Pt(20)
            widths = sorted(
                [
                    min(
                        pptx.util.Inches(font_size_20pt.inches * a_len),
                        width
                    ) for a_len in lengths
                ]
            )
            width = widths[len(widths) // 2]

            for step in steps:
                shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height)
                shape.text = step.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
                top += height + INCHES_0_3
                left += INCHES_0_5

    return True


def _handle_key_message(
        the_slide: pptx.slide.Slide,
        slide_json: dict,
        slide_width_inch: float,
        slide_height_inch: float
):
    """
    Add a shape to display the key message in the slide, if available.

    :param the_slide: The slide to be processed.
    :param slide_json: The content of the slide as JSON data.
    :param slide_width_inch: The width of the slide in inches.
    :param slide_height_inch: The height of the slide in inches.
    """

    if 'key_message' in slide_json and slide_json['key_message']:
        height = pptx.util.Inches(1.6)
        width = pptx.util.Inches(slide_width_inch / 2.3)
        top = pptx.util.Inches(slide_height_inch - height.inches - 0.1)
        left = pptx.util.Inches((slide_width_inch - width.inches) / 2)
        shape = the_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        shape.text = slide_json['key_message']


def _get_slide_width_height_inches(presentation: pptx.Presentation) -> Tuple[float, float]:
    """
    Get the dimensions of a slide in inches.

    :param presentation: The presentation object.
    :return: The width and the height.
    """

    slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
    slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
    # logger.debug('Slide width: %f, height: %f', slide_width_inch, slide_height_inch)

    return slide_width_inch, slide_height_inch


if __name__ == '__main__':
    _JSON_DATA = '''
    {
    "title": "The Fascinating World of Chess",
    "slides": [
        {
            "heading": "Introduction to Chess",
            "bullet_points": [
                "Chess is a strategic board game played between two players.",
                [
                    "Each player begins the game with 16 pieces: one king, one queen, two rooks, two knights, two bishops, and eight pawns.",
                    "The goal of the game is to checkmate your opponent's king. This means the king is in a position to be captured (in 'check') but has no move to escape (mate)."
                ],
                "Chess is believed to have originated in northern India in the 6th century AD."
            ],
            "key_message": "Understanding the basics of chess is crucial before delving into strategies.",
            "img_keywords": "chessboard, chess pieces, king, queen, rook, knight, bishop, pawn"
        },
        {
            "heading": "The Chessboard",
            "bullet_points": [
                "The chessboard is made up of 64 squares in an 8x8 grid.",
                "Each player starts with their pieces on their home rank (row).",
                "The board is divided into two camps: one for each player."
            ],
            "key_message": "Knowing the layout of the chessboard is essential for understanding piece movement.",
            "img_keywords": "chessboard layout, 8x8 grid, home rank, player camps"
        },
        {
            "heading": "Movement of Pieces",
            "bullet_points": [
                ">> Each piece moves differently. Learning these movements is key to playing chess.",
                ">> The king moves one square in any direction.",
                ">> The queen combines the moves of the rook and bishop.",
                ">> The rook moves horizontally or vertically along a rank or file.",
                ">> The bishop moves diagonally.",
                ">> The knight moves in an L-shape: two squares in a horizontal or vertical direction, then one square perpendicular to that.",
                ">> The pawn moves forward one square, but captures diagonally.",
                ">> Pawns have the initial option of moving two squares forward on their first move."
            ],
            "key_message": "Understanding how each piece moves is fundamental to playing chess.",
            "img_keywords": "chess piece movements, king, queen, rook, bishop, knight, pawn"
        },
        {
            "heading": "Special Moves",
            "bullet_points": [
                {
                    "heading": "Castling",
                    "bullet_points": [
                        "Castling is a unique move involving the king and a rook.",
                        "It involves moving the king two squares towards a rook, then moving that rook to the square the king skipped over."
                    ]
                },
                {
                    "heading": "En Passant",
                    "bullet_points": [
                        "En passant is a special pawn capture move.",
                        "It occurs when a pawn moves two squares forward from its starting position and lands beside an opponent's pawn, which could have captured it if the first pawn had only moved one square forward."
                    ]
                }
            ],
            "key_message": "Understanding these special moves can add depth to your chess strategy.",
            "img_keywords": "castling, en passant, special chess moves"
        },
        {
            "heading": "Chess Notation",
            "bullet_points": [
                "Chess notation is a system used to record and communicate chess games.",
                "It uses algebraic notation, where each square on the board is identified by a letter and a number.",
                "Pieces are identified by their initial letters: K for king, Q for queen, R for rook, B for bishop, N for knight, and P for pawn."
            ],
            "key_message": "Learning chess notation is helpful for recording, analyzing, and discussing games.",
            "img_keywords": "chess notation, algebraic notation, chess symbols"
        },
        {
            "heading": "Chess Strategies",
            "bullet_points": [
                "Develop your pieces quickly and efficiently.",
                "Control the center of the board.",
                "Castle early to protect your king.",
                "Keep your king safe.",
                "Think ahead and plan your moves."
            ],
            "key_message": "Following these strategies can help improve your chess skills.",
            "img_keywords": "chess strategies, piece development, center control, king safety, planning ahead"
        },
        {
            "heading": "Chess Tactics",
            "bullet_points": [
                "Fork: attacking two enemy pieces with the same move.",
                "Pin: restricting the movement of an enemy piece.",
                "Skewer: forcing an enemy piece to move away from a threatened piece.",
                "Discovered attack: moving a piece to reveal an attack by another piece behind it."
            ],
            "key_message": "Mastering these tactics can help you gain an advantage in games.",
            "img_keywords": "chess tactics, fork, pin, skewer, discovered attack"
        },
        {
            "heading": "Chess Openings",
            "bullet_points": [
                {
                    "heading": "Italian Game",
                    "bullet_points": [
                        "1. e4 e5",
                        "2. Nf3 Nc6",
                        "3. Bc4 Bc5"
                    ]
                },
                {
                    "heading": "Ruy Lopez",
                    "bullet_points": [
                        "1. e4 e5",
                        "2. Nf3 Nc6",
                        "3. Bb5"
                    ]
                }
            ],
            "key_message": "Learning popular chess openings can help you start games effectively.",
            "img_keywords": "chess openings, Italian Game, Ruy Lopez"
        },
        {
            "heading": "Chess Endgames",
            "bullet_points": [
                {
                    "heading": "King and Pawn Endgame",
                    "bullet_points": [
                        "This endgame involves a king and one or more pawns against a lone king.",
                        "The goal is to promote a pawn to a new queen."
                    ]
                },
                {
                    "heading": "Rook Endgame",
                    "bullet_points": [
                        "This endgame involves a rook against a lone king.",
                        "The goal is to checkmate the opponent's king using the rook."
                    ]
                }
            ],
            "key_message": "Understanding common chess endgames can help you win games.",
            "img_keywords": "chess endgames, king and pawn endgame, rook endgame"
        },
        {
            "heading": "Conclusion",
            "bullet_points": [
                "Chess is a complex game that requires strategy, tactics, and planning.",
                "Understanding the rules, piece movements, and common strategies can help improve your chess skills.",
                "Practice regularly to improve your game."
            ],
            "key_message": "To excel at chess, one must understand its fundamentals and practice regularly.",
            "img_keywords": "chess fundamentals, chess improvement, regular practice"
        }
    ]
}
'''

    temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    path = pathlib.Path(temp.name)

    generate_powerpoint_presentation(
        json5.loads(_JSON_DATA),
        output_file_path=path,
        slides_template='Minimalist Sales Pitch'
    )
    print(f'File path: {path}')

    temp.close()
