from typing import List, Tuple
import json5
import pptx
import re
import yaml

from global_config import GlobalConfig


PATTERN = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)


def remove_slide_number_from_heading(header: str) -> str:
    """
    Remove the slide number from a given slide header.

    :param header: The header of a slide
    """

    if PATTERN.match(header):
        idx = header.find(':')
        header = header[idx + 1:]

    return header


def generate_powerpoint_presentation(
        structured_data: str,
        as_yaml: bool,
        slides_template: str,
        output_file_name: str
) -> List:
    """
    Create and save a PowerPoint presentation file containing the contents in JSON or YAML format.

    :param structured_data: The presentation contents as "JSON" (may contain trailing commas) or YAML
    :param as_yaml: True if the input data is in YAML format; False if it is in JSON format
    :param slides_template: The PPTX template to use
    :param output_file_name: The name of the PPTX file to save as
    :return A list of presentation title and slides headers
    """

    if as_yaml:
        # Avoid YAML mode: nested bullets can lead to incorrect YAML generation
        try:
            parsed_data = yaml.safe_load(structured_data)
        except yaml.parser.ParserError as ype:
            print(f'*** YAML parse error: {ype}')
            parsed_data = {'title': '', 'slides': []}
    else:
        # The structured "JSON" might contain trailing commas, so using json5
        parsed_data = json5.loads(structured_data)

    print(f"*** Using PPTX template: {GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file']}")
    presentation = pptx.Presentation(GlobalConfig.PPTX_TEMPLATE_FILES[slides_template]['file'])

    # The title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = parsed_data['title']
    print(f'Title is: {title.text}')
    subtitle.text = 'by Myself and SlideDeck AI :)'
    all_headers = [title.text, ]

    # background = slide.background
    # background.fill.solid()
    # background.fill.fore_color.rgb = RGBColor.from_string('C0C0C0')  # Silver
    # title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)  # Navy blue

    # Add contents in a loop
    for a_slide in parsed_data['slides']:
        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = remove_slide_number_from_heading(a_slide['heading'])
        all_headers.append(title_shape.text)
        text_frame = body_shape.text_frame

        # The bullet_points may contain a nested hierarchy of JSON arrays
        # In some scenarios, it may contain objects (dictionaries) because the LLM generated so
        #  ^ The second scenario is not covered

        flat_items_list = get_flat_list_of_contents(a_slide['bullet_points'], level=0)

        for an_item in flat_items_list:
            paragraph = text_frame.add_paragraph()
            paragraph.text = an_item[0]
            paragraph.level = an_item[1]

    # The thank-you slide
    last_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(last_slide_layout)
    title = slide.shapes.title
    title.text = 'Thank you!'

    presentation.save(output_file_name)

    return all_headers


def get_flat_list_of_contents(items: list, level: int) -> List[Tuple]:
    """
    Flatten a (hierarchical) list of bullet points to a single list containing each item and its level.

    :param items: A bullet point (string or list)
    :param level: The current level of hierarchy
    :return: A list of (bullet item text, hierarchical level) tuples
    """

    flat_list = []

    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list = flat_list + get_flat_list_of_contents(item, level + 1)

    return flat_list


if __name__ == '__main__':
    # bullets = [
    #     'Description',
    #     'Types',
    #     [
    #         'Type A',
    #         'Type B'
    #     ],
    #     'Grand parent',
    #     [
    #         'Parent',
    #         [
    #             'Grand child'
    #         ]
    #     ]
    # ]

    # output = get_flat_list_of_contents(bullets, level=0)
    # for x in output:
    #     print(x)

    json_data = '''
    {
    "title": "Understanding AI",
    "slides": [
        {
            "heading": "Introduction",
            "bullet_points": [
                "Brief overview of AI",
                [
                    "Importance of understanding AI"
                ]
            ]
        },
        {
            "heading": "What is AI?",
            "bullet_points": [
                "Definition of AI",
                [
                    "Types of AI",
                    [
                        "Narrow or weak AI",
                        "General or strong AI"
                    ]
                ],
                "Differences between AI and machine learning"
            ]
        },
        {
            "heading": "How AI Works",
            "bullet_points": [
                "Overview of AI algorithms",
                [
                    "Types of AI algorithms",
                    [
                        "Rule-based systems",
                        "Decision tree systems",
                        "Neural networks"
                    ]
                ],
                "How AI processes data"
            ]
        },
        {
            "heading": "Pros of AI",
            "bullet_points": [
                "Increased efficiency and productivity",
                "Improved accuracy and precision",
                "Enhanced decision-making capabilities",
                "Personalized experiences"
            ]
        },
        {
            "heading": "Cons of AI",
            "bullet_points": [
                "Job displacement and loss of employment",
                "Bias and discrimination",
                "Privacy and security concerns",
                "Dependence on technology"
            ]
        },
        {
            "heading": "Future Prospects of AI",
            "bullet_points": [
                "Advancements in fields such as healthcare and finance",
                "Increased use"
            ]
        }
    ]
}'''
    generate_powerpoint_presentation(
        json5.loads(json_data),
        as_yaml=False,
        output_file_name='test.pptx',
        slides_template='Blank'
    )
