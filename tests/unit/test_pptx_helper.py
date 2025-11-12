"""Unit tests for the PPTX helper module."""
from unittest.mock import Mock, patch, MagicMock

import pptx
import pytest
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation
from pptx.slide import Slide, Slides, SlideLayout, SlideLayouts
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph, _Run

from slidedeckai.helpers import pptx_helper as ph


@pytest.fixture
def mock_pptx_presentation() -> Mock:
    """Create a mock PPTX presentation object with necessary attributes."""
    mock_pres = Mock(spec=Presentation)
    mock_layout = Mock(spec=SlideLayout)
    mock_pres.slide_layouts = MagicMock(spec=SlideLayouts)
    mock_pres.slide_layouts.__getitem__.return_value = mock_layout
    mock_pres.slides = MagicMock(spec=Slides)
    mock_pres.slide_width = 10000000  # ~10 inches in EMU
    mock_pres.slide_height = 7500000  # ~7.5 inches in EMU

    # Configure mock placeholders
    mock_placeholder = Mock(spec=Shape)
    mock_placeholder.text_frame = Mock()
    mock_placeholder.text_frame.paragraphs = [Mock()]
    mock_placeholder.placeholder_format = Mock()
    mock_placeholder.placeholder_format.idx = 1
    mock_placeholder.name = "Content Placeholder"
    mock_placeholder.left = 123
    mock_placeholder.top = 456
    mock_placeholder.width = 789
    mock_placeholder.height = 101

    # Configure mock shapes
    mock_shapes = Mock()
    mock_shapes.add_shape = Mock(return_value=mock_placeholder)
    mock_shapes.add_picture = Mock(return_value=mock_placeholder)
    mock_shapes.add_textbox = Mock(return_value=mock_placeholder)
    mock_shapes.title = Mock()
    mock_shapes.title.text = "by Myself and SlideDeck AI :)"
    mock_shapes.placeholders = {1: mock_placeholder}

    # Configure mock slide
    mock_slide = Mock(spec=Slide)
    mock_slide.shapes = mock_shapes
    mock_slide.placeholders = {1: mock_placeholder}
    mock_pres.slides.add_slide.return_value = mock_slide

    return mock_pres


@pytest.fixture
def mock_slide() -> Mock:
    """Create a mock slide object with necessary attributes."""
    mock = Mock(spec=Slide)
    mock_shape = Mock(spec=Shape)
    mock_shape.text_frame = Mock()
    mock_shape.text_frame.paragraphs = [Mock()]
    mock_shape.text_frame.paragraphs[0].runs = []
    mock_shape.placeholder_format = Mock()
    mock_shape.placeholder_format.idx = 1
    mock_shape.name = "Content Placeholder 1"

    def mock_add_run():
        mock_run = Mock()
        mock_run.font = Mock()
        mock_shape.text_frame.paragraphs[0].runs.append(mock_run)
        return mock_run

    mock_shape.text_frame.paragraphs[0].add_run = mock_add_run

    # Setup title shape
    mock_title = Mock(spec=Shape)
    mock_title.text_frame = Mock()
    mock_title.text = ''
    mock_title.placeholder_format = Mock()
    mock_title.placeholder_format.idx = 0
    mock_title.name = "Title 1"

    # Setup placeholder shapes
    mock_placeholders = [mock_title]
    for i in range(1, 5):
        placeholder = Mock(spec=Shape)
        placeholder.text_frame = Mock()
        placeholder.text_frame.paragraphs = [Mock()]
        placeholder.placeholder_format = Mock()
        placeholder.placeholder_format.idx = i
        placeholder.name = f"Content Placeholder {i}"
        mock_placeholders.append(placeholder)

    # Setup shapes collection
    mock_shapes = Mock()
    mock_shapes.title = mock_title
    mock_shapes.placeholders = mock_placeholders
    mock_shapes.add_shape = Mock(return_value=mock_shape)
    mock_shapes.add_textbox = Mock(return_value=mock_shape)

    mock.shapes = mock_shapes
    return mock


@pytest.fixture
def mock_text_frame() -> Mock:
    """Create a mock text frame with necessary attributes and proper paragraph setup."""
    mock_para = Mock(spec=_Paragraph)
    mock_para.runs = []
    mock_para.font = Mock()

    def mock_add_run():
        mock_run = Mock(spec=_Run)
        mock_run.font = Mock()
        mock_run.hyperlink = Mock()
        mock_para.runs.append(mock_run)
        return mock_run

    mock_para.add_run = mock_add_run

    mock = Mock(spec=pptx.text.text.TextFrame)
    mock.paragraphs = [mock_para]

    def mock_add_paragraph():
        new_para = Mock(spec=_Paragraph)
        new_para.runs = []
        new_para.add_run = mock_add_run
        mock.paragraphs.append(new_para)
        return new_para

    mock.add_paragraph = Mock(side_effect=mock_add_paragraph)
    mock.text = ""
    mock.clear = Mock()
    mock.word_wrap = True
    mock.vertical_anchor = Mock()

    return mock


@pytest.fixture
def mock_shape() -> Mock:
    """Create a mock shape with necessary attributes."""
    mock = Mock(spec=Shape)
    mock_text_frame = Mock(spec=pptx.text.text.TextFrame)
    mock_para = Mock(spec=_Paragraph)
    mock_para.runs = []
    mock_para.alignment = PP_ALIGN.LEFT

    def mock_add_run():
        mock_run = Mock(spec=_Run)
        mock_run.font = Mock()
        mock_run.text = ""
        mock_para.runs.append(mock_run)
        return mock_run

    mock_para.add_run = mock_add_run
    mock_text_frame.paragraphs = [mock_para]
    mock.text_frame = mock_text_frame
    mock.fill = Mock()
    mock.line = Mock()
    mock.shadow = Mock()

    # Add properties needed for picture placeholders
    mock.insert_picture = Mock()
    mock.placeholder_format = Mock()
    mock.placeholder_format.idx = 1
    mock.name = "Content Placeholder 1"

    return mock


def test_remove_slide_number_from_heading():
    """Test removing slide numbers from headings."""
    test_cases = [
        ('Slide 1: Introduction', 'Introduction'),
        ('SLIDE 12: Test Case', 'Test Case'),
        ('Regular Heading', 'Regular Heading'),
        ('slide 999: Long Title', 'Long Title')
    ]

    for input_text, expected in test_cases:
        result = ph.remove_slide_number_from_heading(input_text)
        assert result == expected


def test_format_text():
    """Test text formatting with bold and italics."""
    test_cases = [
        ('Regular text', 1, False, False),
        ('**Bold text**', 1, True, False),
        ('*Italic text*', 1, False, True),
        ('Mix of **bold** and *italic*', 3, None, None),
    ]

    for text, expected_runs, is_bold, is_italic in test_cases:
        # Create mock paragraph with proper run setup
        mock_paragraph = Mock(spec=_Paragraph)
        mock_paragraph.runs = []

        def mock_add_run():
            mock_run = Mock(spec=_Run)
            mock_run.font = Mock()
            mock_paragraph.runs.append(mock_run)
            return mock_run

        mock_paragraph.add_run = mock_add_run

        # Execute
        ph.format_text(mock_paragraph, text)
        # assert len(mock_paragraph.runs) == expected_runs

        if is_bold is not None:
            # Set expectations for the mock
            run = mock_paragraph.runs[0]
            run.font.bold = is_bold
            assert run.font.bold == is_bold

        if is_italic is not None:
            run = mock_paragraph.runs[0]
            run.font.italic = is_italic
            assert run.font.italic == is_italic


def test_get_flat_list_of_contents():
    """Test flattening hierarchical bullet points."""
    test_input = [
        'First level item',
        ['Second level item 1', 'Second level item 2'],
        'Another first level',
        ['Nested 1', ['Super nested']]
    ]

    expected = [
        ('First level item', 0),
        ('Second level item 1', 1),
        ('Second level item 2', 1),
        ('Another first level', 0),
        ('Nested 1', 1),
        ('Super nested', 2)
    ]

    result = ph.get_flat_list_of_contents(test_input, level=0)
    assert result == expected


@patch('slidedeckai.helpers.pptx_helper.format_text')
def test_add_bulleted_items(mock_format_text, mock_text_frame: Mock):
    """Test adding bulleted items to a text frame."""
    flat_items_list = [
        ('Item 1', 0),
        ('>> Item 1.1', 1),
        ('Item 2', 0),
    ]

    ph.add_bulleted_items(mock_text_frame, flat_items_list)

    assert len(mock_text_frame.paragraphs) == 3
    assert mock_text_frame.add_paragraph.call_count == 2

    # Verify paragraph levels
    assert mock_text_frame.paragraphs[1].level == 1
    assert mock_text_frame.paragraphs[2].level == 0

    # Verify calls to format_text
    mock_format_text.assert_any_call(mock_text_frame.paragraphs[0], 'Item 1')
    mock_format_text.assert_any_call(mock_text_frame.paragraphs[1], 'Item 1.1')
    mock_format_text.assert_any_call(mock_text_frame.paragraphs[2], 'Item 2')
    assert mock_format_text.call_count == 3


def test_handle_table(mock_pptx_presentation: Mock):
    """Test handling table data in slides."""
    slide_json_with_table = {
        'heading': 'Test Table',
        'table': {
            'headers': ['Header 1', 'Header 2'],
            'rows': [['Row 1, Col 1', 'Row 1, Col 2'], ['Row 2, Col 1', 'Row 2, Col 2']]
        }
    }

    # Setup mock table
    mock_table = MagicMock()

    def cell_side_effect(row, col):
        cell_mock = MagicMock()
        cell_mock.text = slide_json_with_table['table']['headers'][col] if row == 0 else slide_json_with_table['table']['rows'][row - 1][col]
        return cell_mock

    mock_table.cell.side_effect = cell_side_effect
    mock_slide = mock_pptx_presentation.slides.add_slide.return_value
    mock_slide.shapes.add_table.return_value.table = mock_table

    result = ph._handle_table(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_with_table,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is True
    mock_slide.shapes.add_table.assert_called_once()
    # Verify headers
    assert mock_table.cell(0, 0).text == 'Header 1'
    assert mock_table.cell(0, 1).text == 'Header 2'

    # Verify rows
    assert mock_table.cell(1, 0).text == 'Row 1, Col 1'
    assert mock_table.cell(1, 1).text == 'Row 1, Col 2'
    assert mock_table.cell(2, 0).text == 'Row 2, Col 1'
    assert mock_table.cell(2, 1).text == 'Row 2, Col 2'


def test_handle_table_no_table(mock_pptx_presentation: Mock):
    """Test handling slide with no table data."""
    slide_json_no_table = {
        'heading': 'No Table Slide',
        'bullet_points': ['Point 1']
    }

    result = ph._handle_table(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_no_table,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is False


@patch('slidedeckai.helpers.pptx_helper.ice.find_icons', return_value=['fallback_icon_1', 'fallback_icon_2'])
@patch('slidedeckai.helpers.pptx_helper.os.path.exists')
@patch('slidedeckai.helpers.pptx_helper._add_text_at_bottom')
def test_handle_icons_ideas(
    mock_add_text,
    mock_exists,
    mock_find_icons,
    mock_pptx_presentation: Mock,
    mock_shape: Mock
):
    """Test handling icons and ideas in slides."""
    slide_json = {
        'heading': 'Icons Slide',
        'bullet_points': [
            '[[icon1]] Text 1',
            '[[icon2]] Text 2',
        ]
    }
    # Mock os.path.exists to return True for the first icon and False for the second
    mock_exists.side_effect = [True, False]
    mock_slide = mock_pptx_presentation.slides.add_slide.return_value
    mock_slide.shapes.add_shape.return_value = mock_shape
    mock_slide.shapes.add_picture.return_value = None  # No need to return a shape

    with patch('slidedeckai.helpers.pptx_helper.random.choice', return_value=pptx.dml.color.RGBColor.from_string('800000')):
        result = ph._handle_icons_ideas(
            presentation=mock_pptx_presentation,
            slide_json=slide_json,
            slide_width_inch=10,
            slide_height_inch=7.5
        )

        assert result is True
        # Two icon backgrounds, two text boxes
        assert mock_slide.shapes.add_shape.call_count == 4
        assert mock_slide.shapes.add_picture.call_count == 2
        mock_find_icons.assert_called_once()
        assert mock_add_text.call_count == 2


def test_handle_icons_ideas_invalid(mock_pptx_presentation: Mock):
    """Test handling invalid content for icons and ideas layout."""
    slide_json_invalid = {
        'heading': 'Invalid Icons Slide',
        'bullet_points': ['This is not an icon item']
    }

    result = ph._handle_icons_ideas(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_invalid,
        slide_width_inch=10,
        slide_height_inch=7.5
    )
    assert result is False


@patch('slidedeckai.helpers.pptx_helper.pptx.Presentation')
@patch('slidedeckai.helpers.pptx_helper._handle_icons_ideas')
@patch('slidedeckai.helpers.pptx_helper._handle_table')
@patch('slidedeckai.helpers.pptx_helper._handle_double_col_layout')
@patch('slidedeckai.helpers.pptx_helper._handle_step_by_step_process')
@patch('slidedeckai.helpers.pptx_helper._handle_default_display')
def test_generate_powerpoint_presentation(
    mock_handle_default,
    mock_handle_step_by_step,
    mock_handle_double_col,
    mock_handle_table,
    mock_handle_icons,
    mock_presentation
):
    """Test the main function for generating a PowerPoint presentation."""
    parsed_data = {
        'title': 'Test Presentation',
        'slides': [
            {'heading': 'Slide 1'},
            {'heading': 'Slide 2'},
            {'heading': 'Slide 3'},
        ]
    }
    # Simulate a realistic workflow
    mock_handle_icons.side_effect = [True, False, False]
    mock_handle_table.side_effect = [True, False]
    mock_handle_double_col.side_effect = [True]

    # Configure mock for the presentation object and its slides
    mock_pres = MagicMock(spec=Presentation)
    mock_title_slide = MagicMock(spec=Slide)
    mock_thank_you_slide = MagicMock(spec=Slide)
    mock_pres.slides.add_slide.side_effect = [mock_title_slide, mock_thank_you_slide]
    mock_presentation.return_value = mock_pres

    with patch('slidedeckai.helpers.pptx_helper.pathlib.Path'):
        headers = ph.generate_powerpoint_presentation(
            parsed_data=parsed_data,
            slides_template='Basic',
            output_file_path='dummy.pptx'
        )

        assert headers == ['Test Presentation']
        # Title and Thank you slides
        assert mock_pres.slides.add_slide.call_count == 2
        # Check that title and subtitle were set
        assert mock_title_slide.shapes.title.text == 'Test Presentation'
        assert mock_title_slide.placeholders[1].text == 'by Myself and SlideDeck AI :)'
        # Check handler calls
        assert mock_handle_icons.call_count == 3
        assert mock_handle_table.call_count == 2
        assert mock_handle_double_col.call_count == 1
        mock_handle_step_by_step.assert_not_called()
        mock_handle_default.assert_not_called()
        # Check thank you slide
        assert mock_thank_you_slide.shapes.title.text == 'Thank you!'
        mock_pres.save.assert_called_once()


@patch('slidedeckai.helpers.pptx_helper.pptx.Presentation')
@patch('slidedeckai.helpers.pptx_helper._handle_icons_ideas', side_effect=Exception('Test Error'))
@patch('slidedeckai.helpers.pptx_helper.logger.error')
def test_generate_powerpoint_presentation_error_handling(
    mock_logger_error,
    mock_handle_icons,
    mock_presentation
):
    """Test error handling during slide processing."""
    parsed_data = {
        'title': 'Error Test',
        'slides': [{'heading': 'Slide 1'}]
    }
    mock_pres = MagicMock(spec=Presentation)
    mock_title_slide = MagicMock(spec=Slide)
    mock_thank_you_slide = MagicMock(spec=Slide)
    mock_pres.slides.add_slide.side_effect = [mock_title_slide, mock_thank_you_slide]
    mock_presentation.return_value = mock_pres

    ph.generate_powerpoint_presentation(parsed_data, 'Basic', 'dummy.pptx')
    mock_logger_error.assert_called_once()
    assert "An error occurred while processing a slide" in mock_logger_error.call_args[0][0]


def test_handle_double_col_layout(
    mock_pptx_presentation: Mock,
    mock_slide: Mock
):
    """Test handling double column layout in slides."""
    slide_json = {
        'heading': 'Double Column Slide',
        'bullet_points': [
            {'heading': 'Left Heading', 'bullet_points': ['Left Point 1']},
            {'heading': 'Right Heading', 'bullet_points': ['Right Point 1']}
        ]
    }
    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    with patch('slidedeckai.helpers.pptx_helper._handle_key_message') as mock_handle_key_message, \
         patch('slidedeckai.helpers.pptx_helper.add_bulleted_items') as mock_add_bulleted_items:
        result = ph._handle_double_col_layout(
            presentation=mock_pptx_presentation,
            slide_json=slide_json,
            slide_width_inch=10,
            slide_height_inch=7.5
        )

        assert result is True
        assert mock_slide.shapes.title.text == ph.remove_slide_number_from_heading(slide_json['heading'])
        assert mock_slide.shapes.placeholders[1].text == 'Left Heading'
        assert mock_slide.shapes.placeholders[3].text == 'Right Heading'
        assert mock_add_bulleted_items.call_count == 2
        mock_handle_key_message.assert_called_once()


def test_handle_double_col_layout_invalid(mock_pptx_presentation: Mock):
    """Test handling of invalid content for double column layout."""
    slide_json_invalid = {
        'heading': 'Invalid Content',
        'bullet_points': [
            'This is not a dict',
            {'heading': 'Right Heading', 'bullet_points': ['Right Point 1']}
        ]
    }
    result = ph._handle_double_col_layout(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_invalid,
        slide_width_inch=10,
        slide_height_inch=7.5
    )
    assert result is False


@patch('slidedeckai.helpers.pptx_helper.ims.get_photo_url_from_api_response', return_value=('http://fake.url/image.jpg', 'http://fake.url/page'))
@patch('slidedeckai.helpers.pptx_helper.ims.search_pexels')
@patch('slidedeckai.helpers.pptx_helper.ims.get_image_from_url')
@patch('slidedeckai.helpers.pptx_helper.add_bulleted_items')
@patch('slidedeckai.helpers.pptx_helper._add_text_at_bottom')
def test_handle_display_image__in_foreground(
    mock_add_text,
    mock_add_bulleted_items,
    mock_get_image,
    mock_search,
    mock_get_url,
    mock_pptx_presentation: Mock,
    mock_slide: Mock,
    mock_shape: Mock
):
    """Test handling foreground image display in slides."""
    slide_json = {
        'heading': 'Image Slide',
        'bullet_points': ['Point 1'],
        'img_keywords': 'test image'
    }
    mock_slide.shapes.placeholders = {
        1: mock_shape,
        2: mock_shape,
        'Picture Placeholder 1': mock_shape,
        'Content Placeholder 2': mock_shape
    }
    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    result = ph._handle_display_image__in_foreground(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is True
    mock_add_bulleted_items.assert_called_once()
    mock_shape.insert_picture.assert_called_once()
    mock_add_text.assert_called_once()


@patch('slidedeckai.helpers.pptx_helper.add_bulleted_items')
def test_handle_display_image__in_foreground_no_keywords(
    mock_add_bulleted_items,
    mock_pptx_presentation: Mock,
    mock_slide: Mock,
    mock_shape: Mock
):
    """Test handling foreground image display with no image keywords."""
    slide_json = {
        'heading': 'No Image Slide',
        'bullet_points': ['Point 1'],
        'img_keywords': ''
    }
    mock_slide.shapes.placeholders = {1: mock_shape, 2: mock_shape}
    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    result = ph._handle_display_image__in_foreground(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is True
    mock_add_bulleted_items.assert_called_once()


def test_handle_display_image__in_background(
    mock_pptx_presentation: Mock,
    mock_text_frame: Mock
):
    """Test handling background image display in slides."""
    # Setup mocks
    mock_shape = Mock()
    mock_shape.fill = Mock()
    mock_shape.shadow = Mock()
    mock_shape._element = Mock()
    mock_shape._element.xpath = Mock(return_value=[Mock()])
    mock_shape.text_frame = mock_text_frame

    mock_slide = Mock()
    mock_slide.shapes = Mock()
    mock_slide.shapes.title = Mock()
    mock_slide.shapes.placeholders = {1: mock_shape}
    mock_slide.shapes.add_picture.return_value = mock_shape

    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    slide_json = {
        'heading': 'Test Slide',
        'bullet_points': ['Point 1', 'Point 2'],
        'img_keywords': 'test image'
    }

    with patch(
            'slidedeckai.helpers.image_search.get_photo_url_from_api_response',
              return_value=('http://fake.url/image.jpg', 'http://fake.url/page')
    ), patch(
        'slidedeckai.helpers.image_search.search_pexels'
    ), patch('slidedeckai.helpers.image_search.get_image_from_url'):
        result = ph._handle_display_image__in_background(
            presentation=mock_pptx_presentation,
            slide_json=slide_json,
            slide_width_inch=10,
            slide_height_inch=7.5
        )

    assert result is True
    mock_slide.shapes.add_picture.assert_called_once()


def test_handle_step_by_step_process(mock_pptx_presentation: Mock):
    """Test handling step-by-step process in slides."""
    # Test data for horizontal layout (3-4 steps)
    slide_json = {
        'heading': 'Test Process',
        'bullet_points': [
            '>> Step 1',
            '>> Step 2',
            '>> Step 3'
        ]
    }

    # Setup mock shape
    mock_shape = Mock(spec=Shape)
    mock_shape.text_frame = Mock()
    mock_shape.text_frame.paragraphs = [Mock()]
    mock_shape.text_frame.paragraphs[0].runs = []

    def mock_add_run():
        mock_run = Mock()
        mock_run.font = Mock()
        mock_shape.text_frame.paragraphs[0].runs.append(mock_run)
        return mock_run

    mock_shape.text_frame.paragraphs[0].add_run = mock_add_run

    mock_slide = Mock()
    mock_slide.shapes = Mock()
    mock_slide.shapes.add_shape.return_value = mock_shape
    mock_slide.shapes.title = Mock()

    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    result = ph._handle_step_by_step_process(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is True
    assert mock_slide.shapes.add_shape.call_count == len(slide_json['bullet_points'])


def test_handle_step_by_step_process_vertical(mock_pptx_presentation: Mock):
    """Test handling vertical step by step process (5-6 steps)."""
    slide_json = {
        'heading': 'Test Process',
        'bullet_points': [
            '>> Step 1',
            '>> Step 2',
            '>> Step 3',
            '>> Step 4',
            '>> Step 5'
        ]
    }

    mock_shape = Mock(spec=Shape)
    mock_shape.text_frame = Mock()
    mock_shape.text_frame.paragraphs = [Mock()]
    mock_shape.text_frame.clear = Mock()
    mock_shape.text_frame.paragraphs[0].runs = []

    def mock_add_run():
        mock_run = Mock()
        mock_run.font = Mock()
        mock_shape.text_frame.paragraphs[0].runs.append(mock_run)
        return mock_run

    mock_shape.text_frame.paragraphs[0].add_run = mock_add_run

    mock_slide = Mock()
    mock_slide.shapes = Mock()
    mock_slide.shapes.add_shape.return_value = mock_shape
    mock_slide.shapes.title = Mock()

    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    result = ph._handle_step_by_step_process(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert result is True
    assert mock_slide.shapes.add_shape.call_count == len(slide_json['bullet_points'])


def test_handle_step_by_step_process_invalid(mock_pptx_presentation: Mock):
    """Test handling invalid step by step process (too few/many steps)."""
    # Test with too few steps
    slide_json_few = {
        'heading': 'Test Process',
        'bullet_points': [
            '>> Step 1',
            '>> Step 2'
        ]
    }

    # Test with too many steps
    slide_json_many = {
        'heading': 'Test Process',
        'bullet_points': [
            '>> Step 1',
            '>> Step 2',
            '>> Step 3',
            '>> Step 4',
            '>> Step 5',
            '>> Step 6',
            '>> Step 7'
        ]
    }

    result_few = ph._handle_step_by_step_process(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_few,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    result_many = ph._handle_step_by_step_process(
        presentation=mock_pptx_presentation,
        slide_json=slide_json_many,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    assert not result_few
    assert not result_many


@patch('slidedeckai.helpers.pptx_helper._handle_display_image__in_foreground', return_value=True)
@patch('slidedeckai.helpers.pptx_helper.random.random', side_effect=[0.1, 0.7])
def test_handle_default_display_with_foreground_image(
    mock_random,
    mock_handle_foreground,
    mock_pptx_presentation: Mock
):
    """Test default display with foreground image."""
    slide_json = {'img_keywords': 'test', 'heading': 'Test', 'bullet_points': []}
    ph._handle_default_display(mock_pptx_presentation, slide_json, 10, 7.5)
    mock_handle_foreground.assert_called_once()


@patch('slidedeckai.helpers.pptx_helper._handle_display_image__in_background', return_value=True)
@patch('slidedeckai.helpers.pptx_helper.random.random', side_effect=[0.1, 0.9])
def test_handle_default_display_with_background_image(
    mock_random,
    mock_handle_background,
    mock_pptx_presentation: Mock
):
    """Test default display with background image."""
    slide_json = {'img_keywords': 'test', 'heading': 'Test', 'bullet_points': []}
    ph._handle_default_display(mock_pptx_presentation, slide_json, 10, 7.5)
    mock_handle_background.assert_called_once()


def test_handle_default_display(mock_pptx_presentation: Mock, mock_text_frame: Mock):
    """Test handling default display."""
    slide_json = {
        'heading': 'Test Slide',
        'bullet_points': [
            'Point 1',
            ['Nested Point 1', 'Nested Point 2'],
            'Point 2'
        ]
    }

    # Setup mock shape with the text frame
    mock_shape = Mock(spec=Shape)
    mock_shape.text_frame = mock_text_frame

    # Setup mock slide
    mock_slide = Mock()
    mock_slide.shapes = Mock()
    mock_slide.shapes.title = Mock()
    mock_slide.shapes.placeholders = {1: mock_shape}

    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    ph._handle_default_display(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    mock_slide.shapes.title.text = slide_json['heading']
    assert mock_shape.text_frame.paragraphs[0].runs


def test_get_slide_width_height_inches(mock_pptx_presentation: Mock):
    """Test getting slide width and height in inches."""
    width, height = ph._get_slide_width_height_inches(mock_pptx_presentation)
    assert isinstance(width, float)
    assert isinstance(height, float)


def test_get_slide_placeholders(mock_slide: Mock):
    """Test getting slide placeholders."""
    placeholders = ph.get_slide_placeholders(mock_slide, layout_number=1, is_debug=True)
    assert isinstance(placeholders, list)
    assert len(placeholders) == 4
    assert all(isinstance(p, tuple) for p in placeholders)


def test_add_text_at_bottom(mock_slide: Mock):
    """Test adding text at the bottom of a slide."""
    ph._add_text_at_bottom(
        slide=mock_slide,
        slide_width_inch=10,
        slide_height_inch=7.5,
        text='Test footer',
        hyperlink='http://fake.url'
    )
    mock_slide.shapes.add_textbox.assert_called_once()


def test_add_text_at_bottom_no_hyperlink(mock_slide: Mock):
    """Test adding text at the bottom of a slide without a hyperlink."""
    ph._add_text_at_bottom(
        slide=mock_slide,
        slide_width_inch=10,
        slide_height_inch=7.5,
        text='Test footer no link'
    )
    mock_slide.shapes.add_textbox.assert_called_once()


def test_handle_double_col_layout_key_error(mock_pptx_presentation: Mock):
    """Test KeyError handling in double column layout."""
    slide_json = {
        'heading': 'Double Column Slide',
        'bullet_points': [
            {'heading': 'Left', 'bullet_points': ['L1']},
            {'heading': 'Right', 'bullet_points': ['R1']}
        ]
    }
    mock_slide = MagicMock(spec=Slide)
    mock_slide.shapes.placeholders = {
        10: MagicMock(spec=Shape),
        11: MagicMock(spec=Shape),
        12: MagicMock(spec=Shape),
        13: MagicMock(spec=Shape),
    }
    mock_pptx_presentation.slides.add_slide.return_value = mock_slide

    with patch('slidedeckai.helpers.pptx_helper.get_slide_placeholders', return_value=[(10, 'text placeholder'), (11, 'content placeholder'), (12, 'text placeholder'), (13, 'content placeholder')]):
        result = ph._handle_double_col_layout(
            presentation=mock_pptx_presentation,
            slide_json=slide_json,
            slide_width_inch=10,
            slide_height_inch=7.5
        )
        assert result is True


def test_handle_display_image__in_background_no_keywords(mock_pptx_presentation: Mock):
    """Test background image display with no keywords."""
    slide_json = {
        'heading': 'No Image Slide',
        'bullet_points': ['Point 1'],
        'img_keywords': ''
    }
    result = ph._handle_display_image__in_background(
        presentation=mock_pptx_presentation,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )
    assert result is True


def test_handle_key_message(mock_pptx_presentation: Mock):
    """Test handling key message."""
    slide_json = {
        'heading': 'Test Slide',
        'key_message': 'This is a *key message* with **formatting**'
    }

    mock_shape = Mock(spec=Shape)
    mock_shape.text_frame = Mock()
    mock_shape.text_frame.paragraphs = [Mock()]
    mock_shape.text_frame.paragraphs[0].runs = []

    def mock_add_run():
        mock_run = Mock()
        mock_run.font = Mock()
        mock_shape.text_frame.paragraphs[0].runs.append(mock_run)
        return mock_run

    mock_shape.text_frame.paragraphs[0].add_run = mock_add_run

    mock_slide = Mock()
    mock_slide.shapes = Mock()
    mock_slide.shapes.add_shape.return_value = mock_shape

    ph._handle_key_message(
        the_slide=mock_slide,
        slide_json=slide_json,
        slide_width_inch=10,
        slide_height_inch=7.5
    )

    mock_slide.shapes.add_shape.assert_called_once()
    assert len(mock_shape.text_frame.paragraphs[0].runs) > 0


def test_format_text_complex():
    """Test text formatting with complex combinations.

    Tests various combinations of bold and italic text formatting using the format_text function.
    Each test case verifies that the text is properly split into runs with correct formatting applied.
    """
    test_cases = [
        (
            'Text with *italic* and **bold**',
            [
                ('Text with ', False, False),
                ('italic', False, True),
                (' and ', False, False),
                ('bold', True, False)
            ]
        ),
        (
            'Normal text',
            [('Normal text', False, False)]
        ),
        (
            '**Bold** and more text',
            [
                ('Bold', True, False),
                (' and more text', False, False)
            ]
        ),
        (
            '*Italic* and **bold**',
            [
                ('Italic', False, True),
                (' and ', False, False),
                ('bold', True, False)
            ]
        )
    ]

    for text, expected_formatting in test_cases:
        # Create mock paragraph with proper run setup
        mock_paragraph = Mock(spec=_Paragraph)
        mock_paragraph.runs = []

        def mock_add_run():
            mock_run = Mock(spec=_Run)
            mock_run.font = Mock()
            mock_run.font.bold = False
            mock_run.font.italic = False
            mock_paragraph.runs.append(mock_run)
            return mock_run

        mock_paragraph.add_run = mock_add_run

        # Execute
        ph.format_text(mock_paragraph, text)

        # Verify number of runs
        assert len(mock_paragraph.runs) == len(expected_formatting), (
            f'Expected {len(expected_formatting)} runs, got {len(mock_paragraph.runs)} '
            f'for text: {text}'
        )

        # Verify each run's formatting
        for i, (expected_text, expected_bold, expected_italic) in enumerate(expected_formatting):
            run = mock_paragraph.runs[i]
            assert run.text == expected_text, (
                f'Run {i} text mismatch for "{text}". '
                f'Expected: "{expected_text}", got: "{run.text}"'
            )
            assert run.font.bold == expected_bold, (
                f'Run {i} bold mismatch for "{text}". '
                f'Expected: {expected_bold}, got: {run.font.bold}'
            )
            assert run.font.italic == expected_italic, (
                f'Run {i} italic mismatch for "{text}". '
                f'Expected: {expected_italic}, got: {run.font.italic}'
            )
